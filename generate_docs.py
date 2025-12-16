
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# Paths
BASE_DIR = r"c:/Users/Rushi/Desktop/GovAssist"
IMG_HOME = r"C:/Users/Rushi/.gemini/antigravity/brain/296156d3-38bb-4668-a5f4-d5b87ca2750a/govassist_home_ui_1765800650369.png"
IMG_CHAT = r"C:/Users/Rushi/.gemini/antigravity/brain/296156d3-38bb-4668-a5f4-d5b87ca2750a/govassist_chat_ui_1765800675763.png"
IMG_DASH = r"C:/Users/Rushi/.gemini/antigravity/brain/296156d3-38bb-4668-a5f4-d5b87ca2750a/govassist_dashboard_ui_1765800699278.png"

def create_ppt():
    try:
        prs = Presentation()
        
        # Read Markdown
        with open(os.path.join(BASE_DIR, "PRESENTATION_SLIDES.md"), "r", encoding="utf-8") as f:
            content = f.read()
        
        slides_content = content.split("## Slide")
        
        for i, slide_text in enumerate(slides_content):
            if not slide_text.strip(): continue
            
            lines = slide_text.strip().split("\n")
            title_line = lines[0].strip() # e.g. "1: Title Slide"
            title_text = title_line.split(":", 1)[1].strip() if ":" in title_line else title_line
            
            # Determine Layout (0=Title, 1=Content)
            # Layout 0 is typically Title Slide (Title + Subtitle)
            # Layout 1 is Title + Content
            
            if "Title Slide" in title_line:
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                if slide.shapes.title:
                    slide.shapes.title.text = "GovAssist AI"
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = "Bridging the Gap Between Citizens and Welfare Schemes\nDecember 2025"
                
                # Add background image if possible, or small logo
                slide.shapes.add_picture(IMG_HOME, Inches(1), Inches(4), width=Inches(8))
                
            else:
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                
                if slide.shapes.title:
                    slide.shapes.title.text = title_text
                    
                # Set Content/Body
                body_text = ""
                for line in lines[1:]:
                    line = line.strip()
                    if line.startswith("**"):
                         body_text += line.replace("**", "") + "\n"
                    elif line.startswith("*") or line.startswith("-"):
                         body_text += line + "\n"
                    elif line and not line.startswith("---"):
                         body_text += line + "\n"
                
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = body_text

                # Inject Images
                if "Our Solution" in title_text or "Key Features" in title_text:
                     # Check which image fits best
                     img = IMG_CHAT if "Solution" in title_text or "Chatbot" in body_text else IMG_DASH
                     # Add picture to the right side
                     try:
                        slide.shapes.add_picture(img, Inches(7), Inches(2), height=Inches(3.5))
                     except Exception as e:
                        print(f"Could not add image to slide: {e}")

        save_path = os.path.join(BASE_DIR, "GovAssist_Presentation.pptx")
        prs.save(save_path)
        print(f"Created PPT: {save_path}")
    except Exception as e:
        print(f"Error creating PPT: {e}")

def create_pdf():
    try:
        save_path = os.path.join(BASE_DIR, "GovAssist_Report.pdf")
        c = canvas.Canvas(save_path, pagesize=letter)
        width, height = letter
        
        # --- TITLE PAGE ---
        c.setFont("Helvetica-Bold", 26)
        c.drawCentredString(width/2, height - 150, "GovAssist AI")
        c.setFont("Helvetica", 16)
        c.drawCentredString(width/2, height - 180, "Project Report & Documentation")
        
        # Hero Image
        try:
            c.drawImage(IMG_HOME, 50, height - 500, width=500, preserveAspectRatio=True, mask='auto', anchor='c')
        except:
            pass
            
        c.setFont("Helvetica", 12)
        c.drawCentredString(width/2, 100, "December 15, 2025")
        c.showPage()
        
        # --- CONTENT ---
        text_object = c.beginText(40, height - 40)
        
        with open(os.path.join(BASE_DIR, "PROJECT_REPORT.md"), "r", encoding="utf-8") as f:
            lines = f.readlines()
            
        for line in lines:
            line = line.strip()
            
            # Check for page space
            if text_object.getY() < 60:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(40, height - 40)
            
            if line.startswith("# "):
                # Finish current page for major sections if needed, or just bold big
                text_object.setFont("Helvetica-Bold", 20)
                text_object.textLine("\n" + line.replace("# ", ""))
                text_object.setFont("Helvetica", 11)
            elif line.startswith("## "):
                text_object.setFont("Helvetica-Bold", 16)
                text_object.textLine("\n" + line.replace("## ", ""))
                text_object.setFont("Helvetica", 11)
            elif line.startswith("### "):
                text_object.setFont("Helvetica-Bold", 12)
                text_object.textLine("\n" + line.replace("### ", ""))
                text_object.setFont("Helvetica", 11)
            else:
                # Basic wrapping
                limit = 95
                if len(line) > limit:
                    chunks = [line[i:i+limit] for i in range(0, len(line), limit)]
                    for chunk in chunks:
                        text_object.textLine(chunk)
                else:
                    text_object.textLine(line)
            
            # Inject Images based on keywords
            if "Intelligent Chatbot" in line and "The core" in line:
                c.drawText(text_object)
                try:
                    c.drawImage(IMG_CHAT, 100, text_object.getY() - 250, width=400, preserveAspectRatio=True, mask='auto')
                except: pass
                text_object.setTextOrigin(40, text_object.getY() - 270)
                if text_object.getY() < 50:
                    c.showPage()
                    text_object = c.beginText(40, height - 40)

            if "Admin Dashboard" in line and "secure admin" in line:
                c.drawText(text_object)
                try:
                    c.drawImage(IMG_DASH, 100, text_object.getY() - 250, width=400, preserveAspectRatio=True, mask='auto')
                except: pass
                text_object.setTextOrigin(40, text_object.getY() - 270)
                if text_object.getY() < 50:
                    c.showPage()
                    text_object = c.beginText(40, height - 40)

        c.drawText(text_object)
        
        # --- REFERENCES ---
        c.showPage()
        c.setFont("Helvetica-Bold", 18)
        c.drawString(40, height - 50, "References")
        c.setFont("Helvetica", 12)
        refs = [
            "1. FastAPI (Backend Framework): https://fastapi.tiangolo.com/",
            "2. React & Vite (Frontend): https://react.dev/",
            "3. OpenAI Whisper (Speech Recognition): https://github.com/openai/whisper",
            "4. Tailwind CSS (Styling): https://tailwindcss.com/",
            "5. SQLAlchemy (Database ORM): https://www.sqlalchemy.org/",
            "6. Framer Motion (Animations): https://www.framer.com/motion/",
            "7. ReportLab (PDF Generation): https://www.reportlab.com/",
            "8. Python-PPTX (Presentation Generation): https://python-pptx.readthedocs.io/"
        ]
        y = height - 80
        for r in refs:
            c.drawString(40, y, r)
            y -= 25
            
        c.save()
        print(f"Created PDF: {save_path}")

    except Exception as e:
        print(f"Error creating PDF: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    create_ppt()
    create_pdf()
